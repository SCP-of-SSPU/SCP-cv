#!/user/bin/env python
# -*- coding: UTF-8 -*-
'''
PPT 解析与转换服务，负责读取 PPT 信息、提取媒体、导出页面图片。
采用方案A：图片底图 + 媒体文件提取 + Qt Multimedia 叠加播放。
@Project : SCP-cv
@File : ppt_processor.py
@Author : Qintsg
@Date : 2026-04-10
'''
from __future__ import annotations

import logging
import subprocess
from pathlib import Path
from typing import Optional

from django.conf import settings

from scp_cv.apps.resources.models import (
    MediaKind,
    ParseState,
    PresentationDocument,
    PresentationPageMedia,
    ResourceFile,
)
from scp_cv.services.executables import get_libreoffice_executable
from scp_cv.services.resource_manager import _get_resource_dir

logger = logging.getLogger(__name__)

# 媒体文件在 PPT 中的 content_type 映射
_VIDEO_CONTENT_TYPES: frozenset[str] = frozenset({
    "video/mp4", "video/avi", "video/x-msvideo", "video/quicktime",
    "video/x-ms-wmv", "video/mpeg", "video/x-matroska",
})
_AUDIO_CONTENT_TYPES: frozenset[str] = frozenset({
    "audio/mpeg", "audio/mp3", "audio/wav", "audio/x-wav",
    "audio/ogg", "audio/x-ms-wma", "audio/aac",
})

# 通过扩展名识别媒体类型的备用映射
_VIDEO_EXTENSIONS: frozenset[str] = frozenset({".mp4", ".avi", ".mov", ".wmv", ".mpeg", ".mkv", ".flv"})
_AUDIO_EXTENSIONS: frozenset[str] = frozenset({".mp3", ".wav", ".ogg", ".wma", ".aac", ".m4a", ".flac"})


class PptProcessorError(Exception):
    """PPT 处理过程中的异常。"""


def parse_and_convert(resource_id: int) -> PresentationDocument:
    """
    完整处理一个 PPT 资源：解析信息 → 导出页面图片 → 提取媒体文件。
    :param resource_id: ResourceFile 主键
    :return: 创建或更新的 PresentationDocument
    """
    try:
        resource = ResourceFile.objects.get(pk=resource_id)
    except ResourceFile.DoesNotExist as not_found:
        raise PptProcessorError(f"资源 id={resource_id} 不存在") from not_found

    resource.parse_state = ParseState.PROCESSING
    resource.save(update_fields=["parse_state"])

    try:
        # 阶段1：用 python-pptx 读取信息和提取媒体
        pptx_info = _extract_pptx_info(resource)

        # 阶段2：用 LibreOffice 导出 PDF，再用 PyMuPDF 转换为逐页 PNG
        derived_path = _convert_to_page_images(resource)

        # 创建或更新 PresentationDocument
        presentation_doc, _created = PresentationDocument.objects.update_or_create(
            resource=resource,
            defaults={
                "total_pages": pptx_info["page_count"],
                "derived_resource_path": str(derived_path),
            },
        )

        # 保存媒体信息
        _save_media_records(presentation_doc, pptx_info["media_items"])

        # 更新资源记录
        resource.page_count = pptx_info["page_count"]
        resource.parse_state = ParseState.READY
        resource.save(update_fields=["page_count", "parse_state"])

        logger.info(
            "PPT 解析完成：「%s」，%d 页，%d 个媒体对象",
            resource.display_name,
            pptx_info["page_count"],
            len(pptx_info["media_items"]),
        )
        return presentation_doc

    except Exception as parse_error:
        resource.parse_state = ParseState.FAILED
        resource.save(update_fields=["parse_state"])
        logger.error("PPT 解析失败（id=%d）：%s", resource_id, parse_error)
        raise PptProcessorError(f"PPT 解析失败：{parse_error}") from parse_error


def _extract_pptx_info(resource: ResourceFile) -> dict[str, object]:
    """
    使用 python-pptx 读取 PPT 的基本信息和页内媒体对象。
    :param resource: 资源文件记录
    :return: 包含 page_count 和 media_items 的字典
    """
    # 延迟导入 pptx 避免模块加载开销
    from pptx import Presentation as PptxPresentation
    from pptx.shapes.base import BaseShape
    from pptx.util import Emu

    original_path = Path(resource.original_path)
    if not original_path.exists():
        raise PptProcessorError(f"原始文件不存在：{original_path}")

    presentation = PptxPresentation(str(original_path))
    page_count = len(presentation.slides)

    # 提取每页中的媒体对象
    media_items: list[dict[str, object]] = []
    resource_dir = _get_resource_dir(resource.pk)
    media_output_dir = resource_dir / "media"
    media_output_dir.mkdir(parents=True, exist_ok=True)

    for slide_index, slide in enumerate(presentation.slides, start=1):
        page_media_order = 0
        for shape in slide.shapes:
            media_info = _detect_shape_media(shape, slide_index, media_output_dir, presentation)
            if media_info is not None:
                media_info["sort_order"] = page_media_order
                media_items.append(media_info)
                page_media_order += 1

    return {
        "page_count": page_count,
        "media_items": media_items,
    }


def _detect_shape_media(
    shape: object,
    page_number: int,
    media_output_dir: Path,
    presentation: object,
) -> Optional[dict[str, object]]:
    """
    检测单个形状是否为媒体对象，若是则提取媒体文件。
    :param shape: pptx 形状对象
    :param page_number: 所在页码
    :param media_output_dir: 媒体文件输出目录
    :param presentation: pptx Presentation 对象
    :return: 媒体信息字典，非媒体返回 None
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 检查是否为媒体形状
    if not hasattr(shape, "shape_type"):
        return None

    shape_type = shape.shape_type
    is_media = shape_type in (MSO_SHAPE_TYPE.MEDIA, MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT)

    # 检查内嵌的媒体元素
    media_element = getattr(shape, "media", None) if hasattr(shape, "media") else None

    if not is_media and media_element is None:
        # 进一步检查是否包含音视频的 relationship
        if not _has_media_relationship(shape):
            return None

    media_name = getattr(shape, "name", f"media_p{page_number}")
    content_type = ""
    media_blob: Optional[bytes] = None
    media_extension = ""

    # 尝试获取内嵌媒体数据
    if media_element is not None:
        content_type = getattr(media_element, "content_type", "")
        media_blob = getattr(media_element, "blob", None)
        media_extension = _guess_extension_from_content_type(content_type)

    # 确定媒体种类
    media_kind = _classify_media_kind(content_type, media_extension, media_name)
    if media_kind is None:
        return None

    # 导出媒体文件
    exported_path = ""
    if media_blob is not None:
        if not media_extension:
            media_extension = ".mp4" if media_kind == MediaKind.VIDEO else ".mp3"
        safe_name = f"p{page_number}_{media_name}{media_extension}"
        export_file_path = media_output_dir / safe_name
        with open(export_file_path, "wb") as media_file:
            media_file.write(media_blob)
        exported_path = str(export_file_path)
        logger.info("提取媒体文件：%s → %s", media_name, export_file_path)

    return {
        "page_number": page_number,
        "media_name": media_name,
        "media_kind": media_kind,
        "media_path": exported_path,
        "sort_order": 0,
    }


def _has_media_relationship(shape: object) -> bool:
    """
    检查形状是否通过 relationship 关联了音视频文件。
    :param shape: pptx 形状对象
    :return: 是否关联媒体
    """
    element = getattr(shape, "_element", None)
    if element is None:
        return False

    # 检查 XML 中是否包含 video/audio 标签
    xml_text = element.xml if hasattr(element, "xml") else ""
    return "a:videoFile" in xml_text or "a:audioFile" in xml_text


def _classify_media_kind(content_type: str, extension: str, name: str) -> Optional[str]:
    """
    根据 content_type、扩展名和名称判断媒体种类。
    :param content_type: MIME 类型
    :param extension: 文件扩展名
    :param name: 媒体名称
    :return: MediaKind 值或 None
    """
    content_type_lower = content_type.lower()
    extension_lower = extension.lower()
    name_lower = name.lower()

    # 优先根据 content_type 判断
    if content_type_lower in _VIDEO_CONTENT_TYPES or content_type_lower.startswith("video/"):
        return MediaKind.VIDEO
    if content_type_lower in _AUDIO_CONTENT_TYPES or content_type_lower.startswith("audio/"):
        return MediaKind.AUDIO

    # 次级通过扩展名判断
    if extension_lower in _VIDEO_EXTENSIONS:
        return MediaKind.VIDEO
    if extension_lower in _AUDIO_EXTENSIONS:
        return MediaKind.AUDIO

    # 最后通过名称推测
    if any(keyword in name_lower for keyword in ("video", "视频", "movie", "film")):
        return MediaKind.VIDEO
    if any(keyword in name_lower for keyword in ("audio", "音频", "sound", "music")):
        return MediaKind.AUDIO

    return None


def _guess_extension_from_content_type(content_type: str) -> str:
    """
    根据 content_type 推测文件扩展名。
    :param content_type: MIME 类型
    :return: 带点号的扩展名
    """
    content_type_map: dict[str, str] = {
        "video/mp4": ".mp4",
        "video/avi": ".avi",
        "video/x-msvideo": ".avi",
        "video/quicktime": ".mov",
        "video/x-ms-wmv": ".wmv",
        "audio/mpeg": ".mp3",
        "audio/mp3": ".mp3",
        "audio/wav": ".wav",
        "audio/x-wav": ".wav",
        "audio/ogg": ".ogg",
    }
    return content_type_map.get(content_type.lower(), "")


def _convert_to_page_images(resource: ResourceFile) -> Path:
    """
    使用 LibreOffice CLI 将 PPT 导出为 PDF，再用 PyMuPDF 逐页转换为 PNG。
    :param resource: 资源文件记录
    :return: 页面图片输出目录路径
    """
    resource_dir = _get_resource_dir(resource.pk)
    pages_output_dir = resource_dir / "pages"
    pages_output_dir.mkdir(parents=True, exist_ok=True)

    original_path = Path(resource.original_path)
    if not original_path.exists():
        raise PptProcessorError(f"原始文件不存在：{original_path}")

    # 阶段1：LibreOffice 导出 PDF
    pdf_path = _export_to_pdf(original_path, resource_dir)

    # 阶段2：PyMuPDF 将 PDF 逐页转换为 PNG
    _pdf_to_page_images(pdf_path, pages_output_dir)

    return pages_output_dir


def _export_to_pdf(ppt_path: Path, output_dir: Path) -> Path:
    """
    调用 LibreOffice CLI 将 PPT 导出为 PDF。
    :param ppt_path: PPT 文件路径
    :param output_dir: 输出目录
    :return: 生成的 PDF 文件路径
    """
    libreoffice_bin = get_libreoffice_executable()
    if libreoffice_bin is None:
        raise PptProcessorError("未检测到 LibreOffice，请先安装并配置 LIBREOFFICE_BIN_PATH")

    # 构造 LibreOffice 命令行
    # --headless: 无界面运行
    # --convert-to pdf: 转换为 PDF
    # --outdir: 输出目录
    command_args = [
        str(libreoffice_bin),
        "--headless",
        "--norestore",
        "--convert-to", "pdf",
        "--outdir", str(output_dir),
        str(ppt_path),
    ]

    logger.info("执行 LibreOffice 转换：%s", " ".join(command_args))

    try:
        completed_process = subprocess.run(
            command_args,
            capture_output=True,
            text=True,
            timeout=120,
            check=False,
        )
    except subprocess.TimeoutExpired as timeout_error:
        raise PptProcessorError("LibreOffice 转换超时（120秒）") from timeout_error

    if completed_process.returncode != 0:
        error_output = completed_process.stderr or completed_process.stdout
        raise PptProcessorError(f"LibreOffice 转换失败（退出码 {completed_process.returncode}）：{error_output}")

    # 查找生成的 PDF 文件
    expected_pdf_name = ppt_path.stem + ".pdf"
    pdf_path = output_dir / expected_pdf_name
    if not pdf_path.exists():
        raise PptProcessorError(f"LibreOffice 转换后未找到 PDF 文件：{pdf_path}")

    logger.info("PDF 导出完成：%s", pdf_path)
    return pdf_path


def _pdf_to_page_images(pdf_path: Path, output_dir: Path) -> int:
    """
    使用 PyMuPDF 将 PDF 逐页转换为 PNG 图片。
    :param pdf_path: PDF 文件路径
    :param output_dir: 页面图片输出目录
    :return: 导出的页数
    """
    import fitz  # PyMuPDF

    pdf_document = fitz.open(str(pdf_path))
    page_count = pdf_document.page_count

    # 使用 2x 缩放获取高分辨率图片，适合大屏显示
    scale_matrix = fitz.Matrix(2.0, 2.0)

    for page_index in range(page_count):
        page = pdf_document[page_index]
        pixmap = page.get_pixmap(matrix=scale_matrix)
        # 页码从 1 开始，文件名格式：page_001.png
        image_file_name = f"page_{page_index + 1:03d}.png"
        image_path = output_dir / image_file_name
        pixmap.save(str(image_path))

    pdf_document.close()
    logger.info("PDF 逐页转换完成：共 %d 页", page_count)
    return page_count


def _save_media_records(
    presentation_doc: PresentationDocument,
    media_items: list[dict[str, object]],
) -> None:
    """
    保存解析出的媒体记录到数据库，先清除旧记录再批量创建。
    :param presentation_doc: PPT 文档模型实例
    :param media_items: 媒体信息字典列表
    """
    # 清除已有的媒体记录（重新解析场景）
    PresentationPageMedia.objects.filter(document=presentation_doc).delete()

    media_objects: list[PresentationPageMedia] = []
    for item in media_items:
        media_objects.append(
            PresentationPageMedia(
                document=presentation_doc,
                page_number=int(item["page_number"]),
                media_name=str(item["media_name"]),
                media_kind=str(item["media_kind"]),
                media_path=str(item["media_path"]),
                sort_order=int(item["sort_order"]),
            )
        )

    if media_objects:
        PresentationPageMedia.objects.bulk_create(media_objects)
        logger.info("保存 %d 条媒体记录", len(media_objects))


def get_page_image_path(resource_id: int, page_number: int) -> Optional[Path]:
    """
    获取指定资源某页的 PNG 图片路径。
    :param resource_id: ResourceFile 主键
    :param page_number: 页码（从 1 开始）
    :return: 图片文件路径，不存在返回 None
    """
    resource_dir = _get_resource_dir(resource_id)
    image_path = resource_dir / "pages" / f"page_{page_number:03d}.png"
    if image_path.exists():
        return image_path
    return None


def get_page_media_list(resource_id: int, page_number: int) -> list[dict[str, object]]:
    """
    获取指定资源指定页的媒体列表。
    :param resource_id: ResourceFile 主键
    :param page_number: 页码
    :return: 媒体信息字典列表
    """
    try:
        resource = ResourceFile.objects.get(pk=resource_id)
    except ResourceFile.DoesNotExist:
        return []

    if not hasattr(resource, "presentation_document"):
        return []

    media_queryset = PresentationPageMedia.objects.filter(
        document=resource.presentation_document,
        page_number=page_number,
    ).order_by("sort_order", "media_name")

    return [
        {
            "id": media_item.pk,
            "media_name": media_item.media_name,
            "media_kind": media_item.media_kind,
            "media_kind_label": media_item.get_media_kind_display(),
            "media_path": media_item.media_path,
            "playback_state": media_item.playback_state,
            "playback_state_label": media_item.get_playback_state_display(),
        }
        for media_item in media_queryset
    ]
