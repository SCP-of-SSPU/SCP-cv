#!/user/bin/env python
# -*- coding: UTF-8 -*-
'''
生成测试用 PPT 文件，用于验证系统的上传、解析、翻页等功能。
运行方式：python tools/testdata/generate_test_ppts.py
@Project : SCP-cv
@File : generate_test_ppts.py
@Author : Qintsg
@Date : 2026-04-10
'''
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 输出目录
OUTPUT_DIR = Path(__file__).resolve().parent


def _add_centered_text(
    slide: object,
    text: str,
    font_size: int = 44,
    color: RGBColor = RGBColor(0x33, 0x33, 0x33),
) -> None:
    """
    在幻灯片中心添加一行文本。
    :param slide: 幻灯片对象
    :param text: 文本内容
    :param font_size: 字号（磅）
    :param color: 字体颜色
    """
    # 使用幻灯片宽度/高度的居中位置
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = True


def _set_slide_background(slide: object, color: RGBColor) -> None:
    """
    设置幻灯片背景色。
    :param slide: 幻灯片对象
    :param color: 背景 RGB 颜色
    """
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def generate_basic_ppt() -> Path:
    """
    生成基础测试 PPT（5 页纯文本），用于验证基本上传和翻页。
    :return: 生成的文件路径
    """
    presentation = Presentation()
    # 设置为 16:9 尺寸
    presentation.slide_width = Emu(12192000)
    presentation.slide_height = Emu(6858000)

    slide_layout = presentation.slide_layouts[6]  # 空白布局
    page_colors = [
        (RGBColor(0xE3, 0xF2, 0xFD), "第 1 页 — 测试封面"),
        (RGBColor(0xE8, 0xF5, 0xE9), "第 2 页 — 内容示例 A"),
        (RGBColor(0xFF, 0xF3, 0xE0), "第 3 页 — 内容示例 B"),
        (RGBColor(0xFC, 0xE4, 0xEC), "第 4 页 — 内容示例 C"),
        (RGBColor(0xF3, 0xE5, 0xF5), "第 5 页 — 结束页"),
    ]

    for background_color, page_text in page_colors:
        slide = presentation.slides.add_slide(slide_layout)
        _set_slide_background(slide, background_color)
        _add_centered_text(slide, page_text)

    output_path = OUTPUT_DIR / "basic_5pages.pptx"
    presentation.save(str(output_path))
    return output_path


def generate_single_page_ppt() -> Path:
    """
    生成单页测试 PPT，用于验证边界情况（翻页到首页/末页）。
    :return: 生成的文件路径
    """
    presentation = Presentation()
    presentation.slide_width = Emu(12192000)
    presentation.slide_height = Emu(6858000)

    slide_layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(slide_layout)
    _set_slide_background(slide, RGBColor(0xE0, 0xF7, 0xFA))
    _add_centered_text(slide, "单页测试 PPT")

    output_path = OUTPUT_DIR / "single_page.pptx"
    presentation.save(str(output_path))
    return output_path


def generate_many_pages_ppt() -> Path:
    """
    生成多页测试 PPT（20 页），用于验证大文件处理和页码跳转。
    :return: 生成的文件路径
    """
    presentation = Presentation()
    presentation.slide_width = Emu(12192000)
    presentation.slide_height = Emu(6858000)

    slide_layout = presentation.slide_layouts[6]

    for page_index in range(1, 21):
        slide = presentation.slides.add_slide(slide_layout)
        # 交替背景色
        hue_value = (page_index * 12) % 256
        background_color = RGBColor(
            min(200 + (page_index * 3) % 56, 255),
            min(200 + (page_index * 7) % 56, 255),
            min(200 + (page_index * 11) % 56, 255),
        )
        _set_slide_background(slide, background_color)
        _add_centered_text(slide, f"第 {page_index} / 20 页")

    output_path = OUTPUT_DIR / "many_pages_20.pptx"
    presentation.save(str(output_path))
    return output_path


def generate_rich_content_ppt() -> Path:
    """
    生成包含多种元素的测试 PPT（标题、正文、形状），用于验证复杂内容渲染。
    :return: 生成的文件路径
    """
    presentation = Presentation()
    presentation.slide_width = Emu(12192000)
    presentation.slide_height = Emu(6858000)

    slide_layout = presentation.slide_layouts[6]

    # 第 1 页：标题页
    slide_title = presentation.slides.add_slide(slide_layout)
    _set_slide_background(slide_title, RGBColor(0x1A, 0x23, 0x7E))
    _add_centered_text(slide_title, "SCP-cv 系统演示", font_size=54, color=RGBColor(0xFF, 0xFF, 0xFF))

    # 添加副标题
    subtitle_box = slide_title.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_run = subtitle_paragraph.add_run()
    subtitle_run.text = "大屏内容播控系统 · 测试演示文稿"
    subtitle_run.font.size = Pt(24)
    subtitle_run.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)

    # 第 2 页：带形状的内容页
    slide_shapes = presentation.slides.add_slide(slide_layout)
    _set_slide_background(slide_shapes, RGBColor(0xFF, 0xFF, 0xFF))
    _add_centered_text(slide_shapes, "形状与色块测试", font_size=36, color=RGBColor(0x21, 0x21, 0x21))

    # 添加彩色矩形
    shape_configs = [
        (Inches(0.5), Inches(5), Inches(2.5), Inches(1), RGBColor(0x42, 0xA5, 0xF5)),
        (Inches(3.5), Inches(5), Inches(2.5), Inches(1), RGBColor(0x66, 0xBB, 0x6A)),
        (Inches(6.5), Inches(5), Inches(2.5), Inches(1), RGBColor(0xFF, 0x72, 0x43)),
    ]
    for rect_left, rect_top, rect_width, rect_height, rect_color in shape_configs:
        shape = slide_shapes.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            rect_left, rect_top, rect_width, rect_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rect_color
        shape.line.fill.background()

    # 第 3 页：长文本页
    slide_text = presentation.slides.add_slide(slide_layout)
    _set_slide_background(slide_text, RGBColor(0xFA, 0xFA, 0xFA))

    long_text_box = slide_text.shapes.add_textbox(Inches(1), Inches(0.8), Inches(8), Inches(5.5))
    long_text_frame = long_text_box.text_frame
    long_text_frame.word_wrap = True

    title_paragraph = long_text_frame.paragraphs[0]
    title_run = title_paragraph.add_run()
    title_run.text = "系统功能概述"
    title_run.font.size = Pt(32)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(0x21, 0x21, 0x21)

    feature_items = [
        "PPT 演示文稿的上传、解析与逐页渲染",
        "SRT 视频流的接收与全屏播放",
        "单屏与左右拼接双屏显示模式",
        "Web 控制台与 gRPC 双通道操控",
        "SSE 实时状态同步",
    ]
    for feature_text in feature_items:
        body_paragraph = long_text_frame.add_paragraph()
        body_paragraph.space_before = Pt(12)
        body_run = body_paragraph.add_run()
        body_run.text = f"• {feature_text}"
        body_run.font.size = Pt(20)
        body_run.font.color.rgb = RGBColor(0x42, 0x42, 0x42)

    output_path = OUTPUT_DIR / "rich_content.pptx"
    presentation.save(str(output_path))
    return output_path


if __name__ == "__main__":
    generated_files = [
        generate_basic_ppt(),
        generate_single_page_ppt(),
        generate_many_pages_ppt(),
        generate_rich_content_ppt(),
    ]
    for generated_path in generated_files:
        print(f"已生成: {generated_path.name} ({generated_path.stat().st_size / 1024:.1f} KB)")
